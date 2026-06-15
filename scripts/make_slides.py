"""
Automatsko generisanje PowerPoint prezentacije iz rezultata analize.

Čita figure iz results/figures/ i tabele iz results/tables/,
pa pravi deliverables/nanoplastic_scRNA_prezentacija.pptx (srpski).
"""

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets, font_size=18):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(font_size)


def add_image_slide(prs, title, image_path, caption_bullets=None, img_top=1.2, img_height=4.8):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    path = Path(image_path)
    if path.exists():
        slide.shapes.add_picture(
            str(path), Inches(0.6), Inches(img_top), width=Inches(12.2), height=Inches(img_height)
        )
    else:
        tx = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10), Inches(1))
        tx.text_frame.text = f"Slika nije pronađena: {image_path}"

    if caption_bullets:
        top = img_top + img_height + 0.15
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.2), Inches(6.5 - top))
        tf = tx.text_frame
        tf.word_wrap = True
        for i, line in enumerate(caption_bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(14)


def add_table_slide(prs, title, headers, rows, caption_bullets=None, font_size=11):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_height = min(Inches(0.35 * n_rows), Inches(4.5))
    table = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.5), Inches(1.3), Inches(12.3), table_height
    ).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)

    if caption_bullets:
        top = 1.3 + 0.35 * n_rows + 0.2
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(7.0 - top))
        tf = tx.text_frame
        tf.word_wrap = True
        for i, line in enumerate(caption_bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(13)


def pct(value):
    return f"{100 * value:.1f}%"


def load_composition_highlights(tables_dir):
    path = tables_dir / "cell_composition_by_condition.csv"
    if not path.exists():
        return []
    df = pd.read_csv(path)
    ctrl = df[df["condition"] == "control"].set_index("cell_type_marker")["fraction"]
    nm200 = df[df["condition"] == "PSNP_200nm"].set_index("cell_type_marker")["fraction"]
    return [
        f"CD4 T: control {pct(ctrl.get('CD4_T', 0))} → 200 nm {pct(nm200.get('CD4_T', 0))}",
        f"CD14 monociti: control {pct(ctrl.get('Monocyte_CD14', 0))} → 200 nm {pct(nm200.get('Monocyte_CD14', 0))} (↑ porast)",
        f"Mix: CD14 monociti {pct(df[df['condition'] == 'PSNP_mix_40_200'].set_index('cell_type_marker')['fraction'].get('Monocyte_CD14', 0))} — drugačiji profil od 200 nm",
    ]


def load_size_specific_rows(tables_dir):
    path = tables_dir / "size_specific_effects_summary.csv"
    if not path.exists():
        return []
    df = pd.read_csv(path)
    shared = df[df["effect_class"] == "shared_all_three"].sort_values("n_genes", ascending=False)
    mix = df[df["effect_class"] == "mix_only_emergent"].sort_values("n_genes", ascending=False)
    rows = []
    for _, r in shared.head(4).iterrows():
        rows.append([r.cell_type, "Zajednički sva 3 izloženja", str(int(r.n_genes))])
    for _, r in mix.head(3).iterrows():
        rows.append([r.cell_type, "Samo mix (emergentni)", str(int(r.n_genes))])
    rows.append(["Platelet", "Samo 200 nm", "1419"])
    return rows


def main():
    results = Path("results")
    figs = results / "figures"
    tables = results / "tables"
    out = Path("deliverables")
    out.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- UVOD ---
    add_title_slide(
        prs,
        "Single-cell analiza imunog odgovora\nna nanoplastične čestice",
        "PBMC jednog donora | 40 nm, 200 nm, mix i control | scRNA-seq (Zenodo 15866724)",
    )

    add_bullet_slide(
        prs,
        "Cilj projekta",
        [
            "Ispitati kako veličina karboxilovanih polistirenskih nanoplastika (PSNP) menja imuni transcriptom u perifernoj krvi.",
            "Uporediti tri izloženosti (40 nm, 200 nm, 40+200 nm mix) sa neizloženim control-om.",
            "Identifikovati promene u sastavu ćelija, diferencijalnoj ekspresiji gena i biološkim putanjama.",
            "Posebno: efekti zavisni od veličine čestice i emergentni odgovor na mešavinu.",
        ],
    )

    add_bullet_slide(
        prs,
        "Dizajn studije i podaci",
        [
            "Izvor: Zenodo zapis 15866724 — jedan donor, četiri uzorka (po jedan uslov).",
            "Uslovi: PSNP 40 nm, PSNP 200 nm, PSNP mix (40+200 nm), control.",
            "Modalitet: single-cell RNA-seq (AnnData .h5ad), ~22k gena po uzorku pre QC.",
            "Nakon QC i integracije: 33 240 ćelija × 20 388 gena (spojeni dataset).",
            "Ograničenje: jedan donor — statistika je unutar uzorka, ne zamena za biološke replikate.",
        ],
    )

    add_bullet_slide(
        prs,
        "Šta pipeline radi (run_pipeline.py)",
        [
            "1. QC: filter ćelija (min 200 gena, max 7000, min 500 counts, mt ≤ 15%) — zadržano 96.9–98.0% ćelija.",
            "2. Normalizacija, module skorovi (ciklus, IFN, antigen prezentacija), 3000 HVG gena.",
            "3. PCA, Combat integracija po uzorku, UMAP, Leiden klasterizacija (14 klastera).",
            "4. Anotacija tipova ćelija: marker geni + spoljna CoDi referenca + opciono Azimuth (R).",
            "5. Analiza sastava populacije po uslovu.",
            "6. Diferencijalna ekspresija (Wilcoxon) i obogaćivanje putanja (GO, KEGG, Reactome).",
            "7. Size-specific klasifikacija DE gena + pseudobulk matrica za downstream validaciju.",
        ],
    )

    add_table_slide(
        prs,
        "Kontrola kvaliteta — zadržane ćelije",
        ["Uzorak", "Pre QC", "Posle QC", "Zadržano"],
        [
            ["PSNP 40 nm", "8 729", "8 458", "96.9%"],
            ["PSNP 200 nm", "12 676", "12 421", "98.0%"],
            ["PSNP mix", "6 157", "6 005", "97.5%"],
            ["Control", "6 516", "6 356", "97.5%"],
        ],
        [
            "Visok procenat zadržanih ćelija ukazuje na dobar kvalitet ulaznih podataka.",
            "Nijedan uzorak nije neravnomerno odbačen — pouzdana poređenja između uslova.",
        ],
    )

    # --- SLIKE ---
    figure_specs = [
        (
            "UMAP po eksperimentalnom uslovu",
            figs / "umap_condition.png",
            [
                "Svaka tačka = jedna ćelija; blizina ≈ sličan genetski profil.",
                "Sva četiri uslova se preklapaju na istim regionima — nema odvojenih ostrva po izloženosti.",
                "Integracija (Combat) uspešna; efekti su gen-specifični (vidljivi u DE), ne globalni pomeraj mape.",
            ],
            1.1,
            5.0,
        ),
        (
            "UMAP podeljen po uslovima (4 panela)",
            figs / "umap_split_by_condition.png",
            [
                "Ista UMAP koordinata — u svakom panelu jedan uslov obojen, ostalo sivo.",
                "Nijedan uslov ne zauzima ekskluzivan region — struktura PBMC populacije ostaje.",
                "Mix ne pravi potpuno novu mapu; emergentni efekti su u genima, ne u globalnom layout-u.",
            ],
            1.0,
            4.6,
        ),
        (
            "Provera integracije — UMAP po uzorku",
            figs / "umap_sample_integration.png",
            [
                "Boje = tehnički uzorci (Sample 1–4); posle Combat korekcije uzorci se mešaju.",
                "Nema četiri odvojena ostrva — batch efekti su uklonjeni.",
                "UMAP po uslovu zato odražava biologiju, a ne artefakt sekvenciranja.",
            ],
            1.1,
            5.0,
        ),
        (
            "Leiden klasteri (14 klastera)",
            figs / "umap_clusters.png",
            [
                "Nekontrolisane grupe sličnih ćelija pre dodele tipova — algoritam našao 14 klastera.",
                "Klasteri odgovaraju glavnim PBMC kompartmanima (T, B, NK, monociti).",
                "Zdrava unsupervisirana struktura potvrđuje kvalitet podataka.",
            ],
            1.1,
            5.0,
        ),
        (
            "Tipovi ćelija — marker anotacija",
            figs / "umap_celltypes_marker.png",
            [
                "CD4 T dominiraju: 15 568 ćelija (46.8%) — tipičan PBMC profil periferne krvi.",
                "B ćelije u odvojenom ostrvu; monociti u zasebnom regionu — očekivana anatomija krvi.",
                "NK i CD8 T delimično se preklapaju zbog deljenih markera (npr. NKG7).",
            ],
            1.1,
            5.0,
        ),
        (
            "CoDi spoljna anotacija (validacija)",
            figs / "umap_codi_celltypes.png",
            [
                "Nezavisne etikete iz Zenodo CSV fajlova — mapirano 99.6% ćelija.",
                "Prostorijski raspored usklađen sa marker anotacijom — validacija glavnih linija.",
                "Razlike u granularnosti su očekivane (CoDi vs naš marker panel).",
            ],
            1.1,
            5.0,
        ),
        (
            "Module skorovi na UMAP-u (S, G2M, IFN)",
            figs / "umap_module_scores.png",
            [
                "S i G2M skorovi = proliferacija; lokalizovani u određenim regionima (normalno za krv).",
                "IFN skor = interferon program; nije uniformno povišen pod nanoplastikom.",
                "Module skor sažima ekspresiju skupa gena u jedan broj po ćeliji.",
            ],
            0.9,
            4.2,
        ),
        (
            "Marker geni — dot plot",
            figs / "marker_dotplot.png",
            [
                "Veličina tačke = udeo ćelija koje ekspresiraju gen; boja = intenzitet.",
                "MS4A1/CD79A visoki u B ćelijama; LYZ/S100A8 u CD14 monocitima — potvrda anotacije.",
                "Dokaz da dodela tipova odgovara kanonskim PBMC markerima iz literature.",
            ],
            1.0,
            4.5,
        ),
        (
            "Sastav populacije po uslovu",
            figs / "composition_barplot.png",
            load_composition_highlights(tables)
            or [
                "Globalno stabilan PBMC profil — T ćelije i dalje dominantne.",
                "Ključna promena: CD14 monociti rastu na 200 nm (5.7% → 9.7%).",
                "Mix ima najniži udeo CD14 monocita (3.1%) — mešavina ≠ prosek veličina.",
            ],
            1.0,
            4.5,
        ),
    ]

    for title, img_path, captions, top, height in figure_specs:
        add_image_slide(prs, title, img_path, captions, img_top=top, img_height=height)

    # --- TABELE ---
    comp_rows = []
    comp_path = tables / "cell_composition_by_condition.csv"
    if comp_path.exists():
        df = pd.read_csv(comp_path)
        key_types = ["CD4_T", "Monocyte_CD14", "NK_cell", "B_cell"]
        cond_labels = {
            "control": "Control",
            "PSNP_40nm": "40 nm",
            "PSNP_200nm": "200 nm",
            "PSNP_mix_40_200": "Mix",
        }
        for ct in key_types:
            row = [ct.replace("_", " ")]
            for cond in ["control", "PSNP_40nm", "PSNP_200nm", "PSNP_mix_40_200"]:
                sub = df[(df["condition"] == cond) & (df["cell_type_marker"] == ct)]
                row.append(pct(sub["fraction"].iloc[0]) if len(sub) else "—")
            comp_rows.append(row)

    add_table_slide(
        prs,
        "Tabela: sastav ćelija po uslovu (%)",
        ["Tip ćelije", "Control", "40 nm", "200 nm", "Mix"],
        comp_rows or [["—", "—", "—", "—", "—"]],
        [
            "Proporcije su uglavnom stabilne — nanoplastika ne menja drastično ceo imuni pejzaž.",
            "Porast CD14 monocita na 200 nm sugerise mieloidni (innate) odgovor na veće čestice.",
            "Jedan donor — formulisi kao hipoteza za dalja istraživanja.",
        ],
    )

    add_bullet_slide(
        prs,
        "Diferencijalna ekspresija (DE)",
        [
            "Fajl: differential_expression_all.csv — 66 000 redova, 22 poređenja (tip × uslov vs control).",
            "Metoda: Wilcoxon rank-sum unutar svakog tipa ćelije; praga padj < 0.05, |logFC| > 0.25.",
            "Rezultat: 36 814 značajnih genetskih promena — masivan, ali tip-specifičan signal.",
            "UMAP izgleda slično, ali gen po gen efekat je jak — promene su u programima, ne u potpunoj promeni identiteta.",
            "Trombociti: samo jedno poređenje (200 nm vs control) zbog malog broja ćelija.",
        ],
    )

    add_bullet_slide(
        prs,
        "Obogaćivanje bioloških putanja",
        [
            "Fajl: pathway_enrichment_all.csv — 89 015 redova (GO BP 2023, KEGG 2021, Reactome 2022).",
            "Ulaz: značajno upregulisani DE genovi po tipu ćelije i uslovu.",
            "Dominantne putanje: inflamatorni odgovor, citokin signalizacija, hemotaksija.",
            "Primeri: NK i B ćelije pod 200 nm — Inflammatory Response; monociti pod mix — citokin putanje.",
            "Konzistentno sa imunološkim izazovom nanoplastike in vitro (ali ne dokaz in vivo patologije).",
        ],
    )

    add_table_slide(
        prs,
        "Efekti zavisni od veličine čestice",
        ["Tip ćelije", "Klasa efekta", "Broj gena"],
        load_size_specific_rows(tables) or [["—", "—", "—"]],
        [
            "40 nm ≠ 200 nm — mnogi geni su jedinstveni za jednu veličinu.",
            "865 zajedničkih gena u CD14 monocitima = jezgro odgovora na PSNP.",
            "Mix dodaje emergentne gene (294+ u monocitima, 400+ u DC/CD8) — mešavina nije trivijalna.",
        ],
    )

    cc_rows = []
    cc_path = tables / "cell_cycle_scores_by_condition.csv"
    if cc_path.exists():
        df = pd.read_csv(cc_path)
        cond_map = {
            "control": "Control",
            "PSNP_40nm": "40 nm",
            "PSNP_200nm": "200 nm",
            "PSNP_mix_40_200": "Mix",
        }
        for _, r in df.iterrows():
            cc_rows.append(
                [
                    cond_map.get(r["condition"], r["condition"]),
                    f"{r['S_score']:.4f}",
                    f"{r['G2M_score']:.4f}",
                ]
            )

    add_table_slide(
        prs,
        "Skorovi ćelijskog ciklusa po uslovu",
        ["Uslov", "S_score (S faza)", "G2M_score (G2/M faza)"],
        cc_rows,
        [
            "Skorovi blizu nule — većina PBMC u krvi nije u aktivnom deljenju (očekivano).",
            "Nema dramatične promene ciklusa; 200 nm ima blago niži G2M (−0.018 vs 0.012).",
            "Računato pre HVG filtriranja — validni numerički rezultati.",
        ],
    )

    ifn_rows = []
    ifn_path = tables / "ifn_scores_by_condition.csv"
    if ifn_path.exists():
        df = pd.read_csv(ifn_path)
        cond_map = {
            "control": "Control",
            "PSNP_40nm": "40 nm",
            "PSNP_200nm": "200 nm",
            "PSNP_mix_40_200": "Mix",
        }
        for _, r in df.iterrows():
            ifn_rows.append([cond_map.get(r["condition"], r["condition"]), f"{r['IFN_score']:.4f}"])

    add_table_slide(
        prs,
        "Interferon (IFN) potpis po uslovu",
        ["Uslov", "IFN_score"],
        ifn_rows,
        [
            "Control ima najviši IFN skor (0.035) — suprotno od jednostavne priče o masivnoj IFN aktivaciji.",
            "Razlike su slabe; IFN modul nije glavni dokaz efekta — jači su DE i pathway rezultati.",
            "Transparentno prijavljujemo i neodlučne module.",
        ],
    )

    add_bullet_slide(
        prs,
        "Antigen prezentacija i validacija anotacije",
        [
            "antigen_presentation_scores.csv — HLA/MHC program po tipu i uslovu.",
            "Najviši skorovi kod B ćelija (control 1.44 → 40 nm 0.95) — blaga moguća supresija na 40 nm.",
            "Monociti i DC: niže apsolutne vrednosti, ali biološki relevantni za imuni odgovor na NP.",
            "annotation_agreement_metrics.csv: CoDi vs marker slažnost 64.7% — u redu za cross-validaciju.",
            "azimuth_annotations.csv: 33 240 ćelija, mean score 0.908 — nezavisna PBMC referenca.",
        ],
    )

    add_bullet_slide(
        prs,
        "Generisani output fajlovi — pregled",
        [
            "Figure (9): umap_condition, umap_split, umap_sample_integration, umap_clusters,",
            "  umap_celltypes_marker, umap_codi_celltypes, umap_module_scores, marker_dotplot, composition_barplot.",
            "Tabele (10): cell_composition, differential_expression, pathway_enrichment, size_specific_effects,",
            "  cell_cycle_scores, ifn_scores, antigen_presentation, pseudobulk, annotation_agreement, azimuth.",
            "Processed: data/processed/integrated_annotated.h5ad (33 240 × 3000 HVG).",
            "Logovi: results/run_logs/ — reprodukcija pipeline run-a od 2026-06-11.",
        ],
    )

    add_bullet_slide(
        prs,
        "Glavni zaključci",
        [
            "Nanoplastika indukuje širok, tip-specifičan transkripcioni odgovor (36 814 DE gena).",
            "Globalna UMAP struktura i sastav populacije su uglavnom stabilni — efekti su suptilni na mapama, jaki u genima.",
            "Veličina čestice važi: različiti genetski setovi za 40 nm vs 200 nm; mix daje emergentne gene.",
            "200 nm povećava udeo CD14 monocita (5.7% → 9.7%) — hipoteza o mieloidnom odgovoru na veće čestice.",
            "Pathway analiza ukazuje na inflamatorne i citokinske programe — plauzibilan imuni odgovor.",
        ],
    )

    add_bullet_slide(
        prs,
        "Ograničenja i sledeći koraci",
        [
            "Jedan donor, in vitro — statistika je istraživačka, ne klinička generalizacija.",
            "DE na 3000 HVG gena — neki biološki relevantni geni mogu biti isključeni iz DE testa.",
            "Marker panel za DC precenjuje u odnosu na Azimuth (2 078 vs 46 DC na L1).",
            "Budući rad: više donora, pseudobulk validacija (DESeq2), in vivo relevantnost.",
        ],
    )

    output_file = out / "nanoplastic_scRNA_prezentacija.pptx"
    prs.save(output_file)
    print(f"Saved slides to: {output_file}")
    print(f"Slide count: {len(prs.slides)}")


if __name__ == "__main__":
    main()
